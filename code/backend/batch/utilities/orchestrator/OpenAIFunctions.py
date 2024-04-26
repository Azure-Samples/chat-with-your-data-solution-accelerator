from typing import List
import json
from pydantic.v1 import BaseModel, Field, validator
import chardet
import mimetypes
from azure.identity import DefaultAzureCredential
from azure.storage.blob import (
    BlobServiceClient,
    ContentSettings,
)

from batch.utilities.helpers import EnvHelper
from .OrchestratorBase import OrchestratorBase
from ..helpers.LLMHelper import LLMHelper
from ..tools.PostPromptTool import PostPromptTool
from ..tools.QuestionAnswerTool import QuestionAnswerTool
from ..tools.TextProcessingTool import TextProcessingTool
from ..tools.ContentSafetyChecker import ContentSafetyChecker
from ..parser.OutputParserTool import OutputParserTool
from ..common.Answer import Answer
from pptx import Presentation
from langchain.agents import tool

env_helper = EnvHelper()
class SlideData(BaseModel):
    title: str = Field(None, description="Title of the slide")
    content: str = Field(None, description="Content of the slide")
    layout: int = Field(0, description="Layout of the slide. Available values: "
                    "0 -> title and subtitle, "
                    "1 -> title and content, "
                    "2 -> section header, "
                    "3 -> two content, "
                    "4 -> Comparison, "
                    "5 -> Title only, "
                    "6 -> Blank, "
                    "7 -> Content with caption, "
                    "8 -> Pic with caption."
                        )
    img_path: str = Field(None, description="Path to the image file for the slide")
    background_path: str = Field(None, description="Path to the background image file for the slide")

    @validator('layout')
    def validate_layout(cls, field):
        if field < 0 or field > 8:
            return ValueError('Layout must be a number from 0 to 8')
        return field
    
    class Config:
        schema_extra = {
            "example": {
                "title": "Slide Title",
                "layout": "0",
                "content": "This is the content of the slide.",
                "img_path": "image.jpg",
                "background_path": "background.jpg"
            }
        }

class PresentationData(BaseModel):
    slides: List[SlideData] = Field(description="List of presentations slides")

    class Config:
        schema_extra = {
            "example": {
                "slides": [
                    {
                        "title": "Slide 1 Title",
                        "content": "Content for slide 1.",
                        "img_path": "slide1_image.jpg",
                        "background_path": "slide1_background.jpg"
                    },
                    {
                        "title": "Slide 2 Title",
                        "content": "Content for slide 2.",
                        "img_path": "slide2_image.jpg",
                        "background_path": "slide2_background.jpg"
                    }
                ]
            }
        }
        
@tool(args_schema=PresentationData)
def create_presentation(slides: List[SlideData]) -> Presentation:
    """Creates PowerPoint presentation"""
    if not slides:
        raise ValueError("Presentation data should have at least one slide")
    
    for index, slide in enumerate(slides):
        if is_slide_data_exceeds_thresholds(slide):
            raise ValueError(f"Slide #{index} exceeded the treshold")

    return try_create_presentation(slides)

def try_create_presentation(slides: List[SlideData]) -> Presentation:
    try:
        prs = Presentation()
        for slide_data in slides:
            layout = int(slide_data.layout)
            slide_layout = prs.slide_layouts[layout]
            slide = prs.slides.add_slide(slide_layout)

            if slide_data.title:
                slide.shapes.title.text = slide_data.title

            if slide_data.content:
                if layout in [1, 3, 4, 7, 8]:
                    content_placeholder = slide.placeholders[1]
                else:
                    content_placeholder = slide.placeholders[0]
                content_placeholder.text = slide_data.content

            #if slide_data.img_path:
                #slide.shapes.add_picture(slide_data.img_path, 0, 0)

            #if slide_data.background_path:
                #slide.background_picture = slide_data.background_path

        return prs
    except Exception as e:
        print("Error:", e)
        return False
    
def is_slide_data_exceeds_thresholds(slide):
    max_characters = 1000  # Define your threshold for maximum characters
    max_images = 3  # Define your threshold for maximum images

    # Calculate the length of text
    text_length = 0 #sum(len(shape.text.strip()) for shape in slide.shapes if shape.has_text_frame)

    # Count the number of images
    num_images = 0 # sum(1 for shape in slide.shapes if shape.shape_type == 13)  # 13 represents Picture shape type

    # Compare with thresholds
    return text_length > max_characters or num_images > max_images

def upload_file(
        bytes_data: bytes, file_name: str
    ):
        if content_type is None:
            content_type = mimetypes.MimeTypes().guess_type(file_name)[0]
            charset = (
                f"; charset={chardet.detect(bytes_data)['encoding']}"
                if content_type == "text/plain"
                else ""
            )
            content_type = content_type if content_type is not None else "text/plain"
        account_name = env_helper.AZURE_BLOB_ACCOUNT_NAME
        account_key = env_helper.AZURE_BLOB_ACCOUNT_KEY
        container_name = env_helper.AZURE_BLOB_CONTAINER_NAME
        if account_name is None or account_key is None or container_name is None:
            raise ValueError(
                "Please provide values for AZURE_BLOB_ACCOUNT_NAME, AZURE_BLOB_ACCOUNT_KEY and AZURE_BLOB_CONTAINER_NAME"
            )
        connect_str = f"DefaultEndpointsProtocol=https;AccountName={account_name};AccountKey={account_key};EndpointSuffix=core.windows.net"
        blob_service_client: BlobServiceClient = (
            BlobServiceClient.from_connection_string(connect_str)
        )
        # Create a blob client using the local file name as the name for the blob
        blob_client = blob_service_client.get_blob_client(
            container=container_name, blob=file_name
        )
        # Upload the created file
        blob_client.upload_blob(
            bytes_data,
            overwrite=True,
            content_settings=ContentSettings(
                content_type=content_type + charset),
        )
            
class OpenAIFunctionsOrchestrator(OrchestratorBase):
    def __init__(self) -> None:
        super().__init__()
        self.content_safety_checker = ContentSafetyChecker()
        self.functions = [
            {
                "name": "search_documents",
                "description": "Retrieve relevant documents to answer user fact-based questions",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "question": {
                            "type": "string",
                            "description": "The user's inquiry, formulated to extract pertinent information from available documents.",
                        },
                        "keywords": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "Relevant keywords, that are list of IT-related terms for precise search",
                        },
                    },
                    "required": ["question", "keywords"],
                },
            },
            {
                "name": "text_processing",
                "description": "Useful when you want to apply a transformation on the text, like translate, summarize, rephrase and so on.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "text": {
                            "type": "string",
                            "description": "The text to be processed",
                        },
                        "operation": {
                            "type": "string",
                            "description": "The operation to be performed on the text. Like Translate to Italian, Summarize, Paraphrase, etc. If a language is specified, return that as part of the operation. Preserve the operation name in the user language.",
                        },
                    },
                    "required": ["text", "operation"],
                },
            },
            {
            "name": "create_presentation",
            "description": "Creates PowerPoint presentation",
            "parameters": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {
                            "type": "string",
                            "description": "Title of the slide"
                        },
                        "content": {
                            "type": "string",
                            "description": "Content of the slide"
                        },
                        "layout": {
                            "type": "integer",
                            "description": "Layout of the slide. Available values: 0 -> title and subtitle, 1 -> title and content, 2 -> section header, 3 -> two content, 4 -> Comparison, 5 -> Title only, 6 -> Blank, 7 -> Content with caption, 8 -> Pic with caption.",
                            "minimum": 0,
                            "maximum": 8
                        },
                        "img_path": {
                            "type": "string",
                            "description": "Path to the image file for the slide"
                        },
                        "background_path": {
                            "type": "string",
                            "description": "Path to the background image file for the slide"
                        }
                    },
                    "required": ["title", "content", "layout"]
                },
                "description": "List of presentations slides"
            }
            }
        ]

    def orchestrate(
        self, user_message: str, chat_history: List[dict], **kwargs: dict
    ) -> dict:
        output_formatter = OutputParserTool()

        # Call Content Safety tool
        if self.config.prompts.enable_content_safety:
            filtered_user_message = (
                self.content_safety_checker.validate_input_and_replace_if_harmful(
                    user_message
                )
            )
            if user_message != filtered_user_message:
                messages = output_formatter.parse(
                    question=user_message,
                    answer=filtered_user_message,
                    source_documents=[],
                )
                return messages

        # Call function to determine route
        llm_helper = LLMHelper()

        system_message = """You help employees to navigate only private information sources, which encompass confidential company documents such as policies, project documentation, technical guides, how-to manuals, and other documentation typical of a large IT company.
        ### IMPORTANT: Your top priority is to utilize the 'search_documents' function with the latest user inquiry for queries concerning these private sources
        ### Instructions for 'search_documents' function:
        1. **Focus on the Most Recent User Inquiry**: Always use the most recent user question as the sole context for the futher steps, we will address to this context as 'user question'. Ignore previous interactions or questions.
        2. **Analyze context**: Carefully read the 'user question' to grasp the intention clearly
        3. **Extract 'question'**:Identify the main intent of the 'user question', keeping it concise and straightforward
            - Ensure the query follows a simple structure suitable for Azure AI Search
            - Optimize the query for effective search results using Azure AI Search best practices
        4. **Extract 'keywords'**:
            - From the 'user question', identify and extract IT-related terms like domains, technologies, frameworks, approaches, testing strategies, etc. without assumptions. If no keywords are available in 'user question', pass an empty array
        
        Call the 'text_processing' function when the user request an operation on the current context, such as translate, summarize, or paraphrase. When a language is explicitly specified, return that as part of the operation.
        When directly replying to the user, always reply in the language the user is speaking.
        """
        # Create conversation history
        messages = [{"role": "system", "content": system_message}]
        for message in chat_history:
            messages.append({"role": message["role"], "content": message["content"]})
        messages.append({"role": "user", "content": user_message})

        result = llm_helper.get_chat_completion_with_functions(
            messages, self.functions, function_call="auto"
        )
        self.log_tokens(
            prompt_tokens=result.usage.prompt_tokens,
            completion_tokens=result.usage.completion_tokens,
        )

        # TODO: call content safety if needed

        if result.choices[0].finish_reason == "function_call":
            if result.choices[0].message.function_call.name == "search_documents":
                func_arguments = json.loads(
                    result.choices[0].message.function_call.arguments
                )
                question = func_arguments["question"]
                # keywords must be a list of strings []
                keywords = func_arguments.get("keywords")

                # run answering chain
                answering_tool = QuestionAnswerTool()
                answer = answering_tool.answer_question(
                    question, chat_history, keywords=keywords
                )

                self.log_tokens(
                    prompt_tokens=answer.prompt_tokens,
                    completion_tokens=answer.completion_tokens,
                )

                # Run post prompt if needed
                if self.config.prompts.enable_post_answering_prompt:
                    post_prompt_tool = PostPromptTool()
                    answer = post_prompt_tool.validate_answer(answer)
                    self.log_tokens(
                        prompt_tokens=answer.prompt_tokens,
                        completion_tokens=answer.completion_tokens,
                    )
                    
            elif result.choices[0].message.function_call.name == "create_presentation":
                func_arguments = json.loads(
                    result.choices[0].message.function_call.arguments
                )
                slides = List[SlideData].parse_raw(func_arguments)
                
                presentation = create_presentation(slides)

                all_texts = [] 
                this_pres_texts = [] 
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            this_pres_texts.append(shape.text)
                all_texts.append(this_pres_texts)
    
                file_name = f'presentation.pptx'
                json_string = json.dumps(all_texts)
                bytes_data = json_string.encode("utf-8")
                upload_file(bytes_data, file_name)
                            
                # run answering chain
                answering_tool = QuestionAnswerTool()
                answer = Answer(
                    question=question,
                    answer="Presentation created",
                    #source_documents=source_documents,
                )
            elif result.choices[0].message.function_call.name == "text_processing":
                text = json.loads(result.choices[0].message.function_call.arguments)[
                    "text"
                ]
                operation = json.loads(
                    result.choices[0].message.function_call.arguments
                )["operation"]
                text_processing_tool = TextProcessingTool()
                answer = text_processing_tool.answer_question(
                    user_message, chat_history, text=text, operation=operation
                )
                self.log_tokens(
                    prompt_tokens=answer.prompt_tokens,
                    completion_tokens=answer.completion_tokens,
                )
        else:
            text = result.choices[0].message.content
            answer = Answer(question=user_message, answer=text)

        
        # Call Content Safety tool
        if self.config.prompts.enable_content_safety:
            filtered_answer = (
                self.content_safety_checker.validate_output_and_replace_if_harmful(
                    answer.answer
                )
            )
            if answer.answer != filtered_answer:
                messages = output_formatter.parse(
                    question=user_message, answer=filtered_answer, source_documents=[]
                )
                return messages

        # Format the output for the UI
        messages = output_formatter.parse(
            question=answer.question,
            answer=answer.answer,
            source_documents=answer.source_documents,
        )
        return messages
