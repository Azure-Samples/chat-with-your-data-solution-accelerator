from datetime import datetime, timedelta
import io
import logging
from pydantic.v1 import BaseModel, Field
from azure.storage.blob import BlobServiceClient, generate_blob_sas, UserDelegationKey

from .EnvHelper import EnvHelper
from pptx import Presentation
from pptx.dml.color import RGBColor
import uuid

env_helper = EnvHelper()
logger = logging.getLogger(__name__)


class ProjectPresentationData(BaseModel):
    name: str = Field(description="Name of the project or slide title.")
    overview: str = Field(
        description="Brief summary of the client's company, including public information."
    )
    challenges: str = Field(
        description="Description of the challenges the company faced and their business goals and needs. This may include issues like recruitment difficulties, market expansion, product failures, or the need for technological expertise."
    )
    technologies: str = Field(
        description="Technologies used in the project, including programming languages, tools, and cloud platforms."
    )
    results: str = Field(
        description="Examples of how Capgemini contributed to the client's business and helped achieve goals, ideally including quantifiable metrics."
    )
    solution: str = Field(
        description="Explanation of how Capgemini addressed the client's needs and assisted them in meeting their goals. This should highlight both business achievements and technical features implemented by our teams."
    )


class PowerPointHelper:

    def __init__(self) -> None:
        self.template_project_presentation_name = "project-template.pptx"
        self.account_name = env_helper.AZURE_BLOB_ACCOUNT_NAME
        self.account_key = env_helper.AZURE_BLOB_ACCOUNT_KEY
        self.container_name = env_helper.AZURE_BLOB_PRESENTATIONS_CONTAINER_NAME

    def load_presentation(self, presentation_name: str) -> Presentation:
        try:
            if env_helper.AZURE_AUTH_TYPE == "rbac":
                raise NotImplementedError(
                    "Presentation retrievement is not implemented for rbac auth type"
                )
            else:
                if (
                    self.account_name is None
                    or self.account_key is None
                    or self.container_name is None
                ):
                    raise ValueError(
                        "Please provide values for AZURE_BLOB_ACCOUNT_NAME, AZURE_BLOB_ACCOUNT_KEY and AZURE_BLOB_PRESENTATIONS_CONTAINER_NAME"
                    )
                connect_str = f"DefaultEndpointsProtocol=https;AccountName={self.account_name};AccountKey={self.account_key};EndpointSuffix=core.windows.net"
                blob_service_client: BlobServiceClient = (
                    BlobServiceClient.from_connection_string(connect_str)
                )

                blob_client = blob_service_client.get_blob_client(
                    container=self.container_name, blob=presentation_name
                )

                blob_data = blob_client.download_blob()
                blob_bytes = blob_data.readall()
                blob_file = io.BytesIO(blob_bytes)
                return Presentation(blob_file)

        except Exception as e:
            logging.error(
                f"Error while retrieving presentation {presentation_name}: {e}"
            )
            raise e

    def save_presentation(self, presentation: Presentation, presentation_name) -> str:
        try:
            if env_helper.AZURE_AUTH_TYPE == "rbac":
                raise NotImplementedError(
                    "Presentation retrievement is not implemented for rbac auth type"
                )
            else:
                if (
                    self.account_name is None
                    or self.account_key is None
                    or self.container_name is None
                ):
                    raise ValueError(
                        "Please provide values for AZURE_BLOB_ACCOUNT_NAME, AZURE_BLOB_ACCOUNT_KEY and AZURE_BLOB_PRESENTATIONS_CONTAINER_NAME"
                    )
                connect_str = f"DefaultEndpointsProtocol=https;AccountName={self.account_name};AccountKey={self.account_key};EndpointSuffix=core.windows.net"
                blob_service_client: BlobServiceClient = (
                    BlobServiceClient.from_connection_string(connect_str)
                )

                blob_client = blob_service_client.get_blob_client(
                    container=self.container_name, blob=presentation_name
                )

                pptx_file = io.BytesIO()
                presentation.save(pptx_file)
                pptx_file.seek(0)

                blob_client.upload_blob(pptx_file)
                blob_sas = generate_blob_sas(
                    self.account_name,
                    self.container_name,
                    presentation_name,
                    account_key=self.account_key,
                    permission="r",
                    expiry=datetime.utcnow() + timedelta(hours=3),
                )

                return (
                    f"https://{self.account_name}.blob.core.windows.net/{self.container_name}/{presentation_name}"
                    + "?"
                    + blob_sas
                )

        except Exception as e:
            logging.error(f"Error while saving presentation {presentation_name}: {e}")
            raise e

    def create_project_presentation(self, projectData: ProjectPresentationData) -> str:
        template = self.load_presentation(
            presentation_name=self.template_project_presentation_name
        )
        self.replace_placeholders_with_data(
            presentation=template, projectData=projectData
        )
        guid = uuid.uuid4()
        presentation_name = f"{projectData.name}{str(guid)}.pptx"
        return self.save_presentation(
            presentation=template, presentation_name=presentation_name
        )

    def replace_placeholders_with_data(
        self, presentation: Presentation, projectData: ProjectPresentationData
    ):
        for slide in presentation.slides:
            for shape in slide.shapes:
                match shape.shape_id:
                    case 31:
                        replaceText(shape.text_frame.paragraphs[0], projectData.name)
                    # case 5:
                    #    shape.text = "project logo"
                    case 24:
                        replaceText(
                            shape.text_frame.paragraphs[1], projectData.overview
                        )
                    case 2:
                        replaceText(
                            shape.text_frame.paragraphs[1], projectData.challenges
                        )
                    case 23:
                        replaceText(
                            shape.text_frame.paragraphs[1], projectData.technologies
                        )
                    case 20:
                        replaceText(shape.text_frame.paragraphs[1], projectData.results)
                    case 18:
                        replaceText(
                            shape.text_frame.paragraphs[1], projectData.solution
                        )


def replaceText(paragraph, source_text):
    for run in paragraph.runs:
        font = run.font
        break
    paragraph.text = source_text
    paragraph.font.name = font.name
    paragraph.font.size = font.size
    paragraph.font.bold = font.bold
    paragraph.font.color.theme_color = font.color.theme_color
    paragraph.font.color.brightness = font.color.brightness
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
